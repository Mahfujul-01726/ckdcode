"""
Generate PowerPoint Presentation from Markdown Slides
For CKD Research Paper
"""

import re
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    import pptx
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call(['pip', 'install', 'python-pptx'])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

def create_presentation():
    """Create PowerPoint presentation from markdown slides"""
    
    # Read markdown slides
    with open('h:\\ckdcode\\presentation_slides.md', 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split into individual slides
    slides_raw = content.split('---\n\n## Slide')
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Chronic Kidney Disease (CKD) Prediction"
    subtitle.text = "Machine Learning-Based Clinical Diagnosis System\n\nA Comparative Study Using Patients Clinical Records\n\nAuthors: Hassan et al. (2023)\n\nHuman-Centric Intelligent Systems 3:92–104"
    
    # Process each slide
    for i, slide_content in enumerate(slides_raw[1:], 1):
        # Parse slide title and content
        lines = slide_content.strip().split('\n')
        
        # Find title (first line after ## Slide X:)
        title_match = re.match(r'(\d+):\s*(.+)', lines[0])
        if not title_match:
            continue
            
        slide_num = title_match.group(1)
        slide_title = title_match.group(2).strip()
        
        # Get content (everything after title)
        content_lines = lines[1:]
        content_text = '\n'.join(content_lines).strip()
        
        # Add slide
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Set content
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()  # Remove default text
        
        # Process content
        process_content(content_text, text_frame)
        
        # Add slide number
        add_slide_number(slide, slide_num)
    
    # Save presentation
    output_path = 'h:\\ckdcode\\CKD_Prediction_Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    
    return output_path

def process_content(content_text, text_frame):
    """Process markdown content and add to text frame"""
    
    lines = content_text.split('\n')
    
    for line in lines:
        line = line.strip()
        
        if not line:
            # Empty line - add paragraph break
            p = text_frame.add_paragraph()
            p.text = ""
            continue
        
        # Check for bullet points
        if line.startswith('✓') or line.startswith('•') or line.startswith('-'):
            p = text_frame.add_paragraph()
            text = line.lstrip('✓•- ').strip()
            
            # Handle bold text
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            
            p.text = text
            p.level = 0
            p.font.size = Pt(18)
            p.font.bold = False
            
        elif line.startswith('|') and '---' in content_text:
            # Table - simplified handling
            p = text_frame.add_paragraph()
            p.text = "[Table/Chart - See visualization files]"
            p.font.size = Pt(14)
            p.font.italic = True
            
        elif '```' in line:
            # Code block - skip
            continue
            
        elif line.startswith('**') or line.endswith('**'):
            # Bold text
            p = text_frame.add_paragraph()
            text = line.strip('*').strip()
            p.text = text
            p.font.size = Pt(20)
            p.font.bold = True
            
        elif '---' in line:
            # Separator
            p = text_frame.add_paragraph()
            p.text = ""
            
        else:
            # Regular text
            p = text_frame.add_paragraph()
            
            # Handle inline formatting
            text = line
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
            text = re.sub(r'`(.+?)`', r'\1', text)
            
            p.text = text
            p.font.size = Pt(18)
            p.font.bold = False

def add_slide_number(slide, number):
    """Add slide number to corner"""
    # Create textbox for slide number
    left = Inches(12.5)
    top = Inches(7.0)
    width = Inches(0.5)
    height = Inches(0.3)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.RIGHT

def create_simple_markdown_presentation():
    """Create a simple markdown-based presentation viewer"""
    
    # Read the presentation markdown
    with open('h:\\ckdcode\\presentation_slides.md', 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split into slides
    slides = content.split('\n---\n\n')
    
    # Create HTML presentation
    html = """<!DOCTYPE html>
<html>
<head>
    <title>CKD Prediction Presentation</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 0; background: #f0f0f0; }
        .slide { display: none; padding: 50px; min-height: 100vh; background: white; box-shadow: 0 0 10px rgba(0,0,0,0.1); }
        .slide.active { display: block; }
        .slide h1 { color: #2c3e50; border-bottom: 3px solid #3498db; padding-bottom: 10px; }
        .slide h2 { color: #34495e; }
        .slide ul { line-height: 1.8; }
        .slide li { margin: 10px 0; }
        .navigation { position: fixed; bottom: 20px; right: 20px; z-index: 1000; }
        .nav-btn { background: #3498db; color: white; border: none; padding: 10px 20px; margin: 0 5px; cursor: pointer; border-radius: 5px; font-size: 16px; }
        .nav-btn:hover { background: #2980b9; }
        .slide-number { position: fixed; bottom: 20px; left: 20px; font-size: 14px; color: #7f8c8d; }
        table { border-collapse: collapse; width: 100%; margin: 20px 0; }
        th, td { border: 1px solid #ddd; padding: 12px; text-align: left; }
        th { background: #3498db; color: white; }
        tr:nth-child(even) { background: #f2f2f2; }
        .highlight { background: #fff3cd; padding: 15px; border-left: 4px solid #ffc107; margin: 20px 0; }
    </style>
</head>
<body>
"""
    
    for i, slide in enumerate(slides):
        # Parse slide
        lines = slide.strip().split('\n')
        
        # Extract title
        title = "Slide " + str(i+1)
        content_lines = []
        
        for line in lines:
            if line.startswith('# '):
                title = line[2:].strip()
            elif line.startswith('## '):
                content_lines.append(f'<h2>{line[3:].strip()}</h2>')
            elif line.startswith('**') and line.endswith('**'):
                content_lines.append(f'<p><strong>{line.strip("*")}</strong></p>')
            elif line.startswith('✓'):
                content_lines.append(f'<li>{line[1:].strip()}</li>')
            elif line.startswith('|') and i < len(lines) - 1 and '|' in lines[min(i+3, len(lines)-1)]:
                # Table
                content_lines.append('<table>')
                for j in range(i, min(i+10, len(lines))):
                    if '|' in lines[j]:
                        content_lines.append('<tr>')
                        for cell in lines[j].split('|')[1:-1]:
                            is_header = False
                            if j+1 < len(lines):
                                is_header = '---' not in lines[j+1]
                            tag = 'th' if is_header else 'td'
                            content_lines.append(f'<{tag}>{cell.strip()}</{tag}>')
                        content_lines.append('</tr>')
                content_lines.append('</table>')
            elif line and not line.startswith('---'):
                content_lines.append(f'<p>{line}</p>')
        
        active = 'active' if i == 0 else ''
        html += f'<div class="slide {active}" id="slide{i}">\n'
        html += f'<h1>{title}</h1>\n'
        html += '\n'.join(content_lines)
        html += '\n</div>\n'
    
    html += """
    <div class="slide-number" id="slideNum">1 / SLIDES_COUNT</div>
    <div class="navigation">
        <button class="nav-btn" onclick="prevSlide()">← Previous</button>
        <button class="nav-btn" onclick="nextSlide()">Next →</button>
    </div>
    <script>
        let currentSlide = 0;
        const slides = document.querySelectorAll('.slide');
        const totalSlides = slides.length;
        
        document.getElementById('slideNum').textContent = `1 / ${totalSlides}`;
        
        function showSlide(n) {
            slides[currentSlide].classList.remove('active');
            currentSlide = (n + totalSlides) % totalSlides;
            slides[currentSlide].classList.add('active');
            document.getElementById('slideNum').textContent = `${currentSlide + 1} / ${totalSlides}`;
        }
        
        function nextSlide() {
            showSlide(currentSlide + 1);
        }
        
        function prevSlide() {
            showSlide(currentSlide - 1);
        }
        
        document.addEventListener('keydown', function(e) {
            if (e.key === 'ArrowRight' || e.key === ' ') nextSlide();
            if (e.key === 'ArrowLeft') prevSlide();
        });
    </script>
</body>
</html>
"""
    
    html = html.replace('SLIDES_COUNT', str(len(slides)))
    
    with open('h:\\ckdcode\\presentation.html', 'w', encoding='utf-8') as f:
        f.write(html)
    
    print("HTML presentation saved to: h:\\ckdcode\\presentation.html")

if __name__ == '__main__':
    print("Generating presentation files...")
    print("\n1. Creating HTML presentation...")
    create_simple_markdown_presentation()
    print("\n2. Creating PowerPoint presentation...")
    try:
        create_presentation()
    except Exception as e:
        print(f"Note: PowerPoint generation requires python-pptx: {e}")
        print("HTML presentation is available as alternative.")
    print("\nDone! Check the output files.")
